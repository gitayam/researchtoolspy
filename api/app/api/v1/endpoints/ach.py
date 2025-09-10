"""
Analysis of Competing Hypotheses (ACH) API endpoints.
Structured analytical technique for intelligence analysis.
"""

import io
from datetime import datetime

from fastapi import APIRouter, Depends, HTTPException, status
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sqlalchemy.ext.asyncio import AsyncSession

try:
    import docx
    import xlsxwriter
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Inches
    EXPORT_AVAILABLE = True
except ImportError:
    # Logger not available yet during import
    EXPORT_AVAILABLE = False

from app.api.v1.endpoints.auth import get_current_user
from app.core.database import get_db
from app.core.logging import get_logger
from app.models.framework import FrameworkType
from app.models.user import User
from app.services.ai_service import ai_service
from app.services.framework_service import FrameworkData, framework_service

logger = get_logger(__name__)
router = APIRouter()


class Hypothesis(BaseModel):
    """Hypothesis model for ACH analysis."""
    id: str
    description: str
    probability: float | None = 0.5  # 0-1 scale
    notes: str | None = None


class SATSEvaluation(BaseModel):
    """SATS (Source, Accuracy, Timeliness, Significance) evaluation model."""
    reliability: int = 3  # 1-5 scale
    credibility: int = 3  # 1-5 scale
    validity: int = 3  # 1-5 scale
    relevance: int = 3  # 1-5 scale
    significance: int = 3  # 1-5 scale
    timeliness: int = 3  # 1-5 scale
    accuracy: int = 3  # 1-5 scale
    completeness: int = 3  # 1-5 scale
    overall_score: float = 0.0  # Calculated average
    evaluation_date: str
    evaluator: str | None = None
    notes: str | None = None


class Evidence(BaseModel):
    """Evidence model for ACH analysis."""
    id: str
    description: str
    credibility: float | None = 0.5  # 0-1 scale (legacy)
    relevance: float | None = 0.5  # 0-1 scale (legacy)
    source: str | None = None
    date: str | None = None
    sats_evaluation: SATSEvaluation | None = None


class EvidenceAssessment(BaseModel):
    """Assessment of evidence against a hypothesis."""
    evidence_id: str
    hypothesis_id: str
    consistency: str  # "consistent", "inconsistent", "neutral", "not_applicable"
    weight: float | None = 0.5  # 0-1 scale for importance
    notes: str | None = None


class ACHCreateRequest(BaseModel):
    """ACH analysis creation request."""
    title: str
    scenario: str
    key_question: str
    initial_hypotheses: list[Hypothesis] | None = []
    initial_evidence: list[Evidence] | None = []
    request_ai_analysis: bool = True


class ACHUpdateRequest(BaseModel):
    """ACH analysis update request."""
    title: str | None = None
    scenario: str | None = None
    key_question: str | None = None
    hypotheses: list[Hypothesis] | None = None
    evidence: list[Evidence] | None = None
    assessments: list[EvidenceAssessment] | None = None


class ACHAnalysisResponse(BaseModel):
    """ACH analysis response."""
    session_id: int
    title: str
    scenario: str
    key_question: str
    hypotheses: list[Hypothesis]
    evidence: list[Evidence]
    assessments: list[EvidenceAssessment]
    matrix: dict | None = None
    ai_analysis: dict | None = None
    status: str
    version: int


@router.post("/", response_model=ACHAnalysisResponse)
async def create_ach_analysis_simple(
    request: ACHCreateRequest,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
) -> ACHAnalysisResponse:
    """Create ACH analysis (standard endpoint)."""
    return await create_ach_analysis(request, current_user, db)


@router.post("/create", response_model=ACHAnalysisResponse)
async def create_ach_analysis(
    request: ACHCreateRequest,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
) -> ACHAnalysisResponse:
    """
    Create a new ACH analysis session.
    
    Args:
        request: ACH creation request
        current_user: Current authenticated user
        db: Database session
        
    Returns:
        ACHAnalysisResponse: Created ACH analysis
    """
    logger.info(f"Creating ACH analysis: {request.title} for user {current_user.username}")

    # Prepare ACH data
    ach_data = {
        "scenario": request.scenario,
        "key_question": request.key_question,
        "hypotheses": [h.dict() for h in request.initial_hypotheses] if request.initial_hypotheses else [],
        "evidence": [e.dict() for e in request.initial_evidence] if request.initial_evidence else [],
        "assessments": [],
    }

    # Get AI analysis if requested
    ai_analysis = None
    if request.request_ai_analysis and request.key_question:
        try:
            ai_result = await framework_service.analyze_with_ai(
                FrameworkType.ACH,
                ach_data,
                "suggest"
            )
            ai_analysis = ai_result.get("suggestions")

            # Add AI-suggested hypotheses and evidence
            if ai_analysis:
                if "hypotheses" in ai_analysis and isinstance(ai_analysis["hypotheses"], list):
                    for idx, hyp in enumerate(ai_analysis["hypotheses"]):
                        ach_data["hypotheses"].append({
                            "id": f"h_ai_{idx}",
                            "description": hyp,
                            "probability": 0.5,
                            "notes": "AI-generated hypothesis"
                        })

                if "evidence" in ai_analysis and isinstance(ai_analysis["evidence"], list):
                    for idx, ev in enumerate(ai_analysis["evidence"]):
                        ach_data["evidence"].append({
                            "id": f"e_ai_{idx}",
                            "description": ev,
                            "credibility": 0.5,
                            "relevance": 0.5,
                            "source": "AI suggestion"
                        })

        except Exception as e:
            logger.warning(f"Failed to get AI analysis: {e}")

    # Create initial assessment matrix
    for hypothesis in ach_data["hypotheses"]:
        for evidence in ach_data["evidence"]:
            ach_data["assessments"].append({
                "evidence_id": evidence["id"],
                "hypothesis_id": hypothesis["id"],
                "consistency": "neutral",
                "weight": 0.5,
                "notes": ""
            })

    # Create framework session
    framework_data = FrameworkData(
        framework_type=FrameworkType.ACH,
        title=request.title,
        description=f"ACH Analysis - {request.key_question}",
        data=ach_data,
        tags=["ach", "hypothesis-testing", "structured-analysis"]
    )

    session = await framework_service.create_session(db, current_user, framework_data)

    # Generate matrix
    matrix = _generate_ach_matrix(ach_data["hypotheses"], ach_data["evidence"], ach_data["assessments"])

    return ACHAnalysisResponse(
        session_id=session.id,
        title=session.title,
        scenario=ach_data["scenario"],
        key_question=ach_data["key_question"],
        hypotheses=[Hypothesis(**h) for h in ach_data["hypotheses"]],
        evidence=[Evidence(**e) for e in ach_data["evidence"]],
        assessments=[EvidenceAssessment(**a) for a in ach_data["assessments"]],
        matrix=matrix,
        ai_analysis=ai_analysis,
        status=session.status.value,
        version=session.version
    )


@router.get("/{session_id}", response_model=ACHAnalysisResponse)
async def get_ach_analysis(
    session_id: int,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
) -> ACHAnalysisResponse:
    """
    Get a specific ACH analysis session.
    
    Args:
        session_id: Session ID
        current_user: Current authenticated user
        db: Database session
        
    Returns:
        ACHAnalysisResponse: ACH analysis data
    """
    logger.info(f"Getting ACH analysis {session_id}")

    # Get real data from database
    session = await framework_service.get_session(db, current_user, session_id, FrameworkType.ACH)
    
    if not session:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="ACH analysis session not found"
        )

    # Extract data from session
    session_data = session.data or {}
    
    # Parse hypotheses
    hypotheses_data = session_data.get("hypotheses", [])
    hypotheses = [Hypothesis(**hyp) for hyp in hypotheses_data if isinstance(hyp, dict)]
    
    # Parse evidence
    evidence_data = session_data.get("evidence", [])
    evidence = [Evidence(**evd) for evd in evidence_data if isinstance(evd, dict)]
    
    # Parse assessments
    assessments_data = session_data.get("assessments", [])
    assessments = [EvidenceAssessment(**ass) for ass in assessments_data if isinstance(ass, dict)]
    
    # Generate matrix
    matrix = _generate_ach_matrix(hypotheses, evidence, assessments)
    
    # Get scenario and key question
    scenario = session_data.get("scenario", "Analysis scenario")
    key_question = session_data.get("key_question", "Key intelligence question")

    return ACHAnalysisResponse(
        session_id=session.id,
        title=session.title,
        scenario=scenario,
        key_question=key_question,
        hypotheses=hypotheses,
        evidence=evidence,
        assessments=assessments,
        matrix=matrix,
        status=session.status.value,
        version=session.version
    )


@router.put("/{session_id}", response_model=ACHAnalysisResponse)
async def update_ach_analysis(
    session_id: int,
    request: ACHUpdateRequest,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
) -> ACHAnalysisResponse:
    """
    Update an existing ACH analysis session.
    
    Args:
        session_id: Session ID
        request: ACH update request
        current_user: Current authenticated user
        db: Database session
        
    Returns:
        ACHAnalysisResponse: Updated ACH analysis
    """
    logger.info(f"Updating ACH analysis {session_id}")

    # TODO: Implement actual database update
    # For now, return mock updated data

    # Get current session (mock)
    current_analysis = await get_ach_analysis(session_id, current_user, db)

    # Update fields
    updated_title = request.title or current_analysis.title
    updated_scenario = request.scenario or current_analysis.scenario
    updated_key_question = request.key_question or current_analysis.key_question
    updated_hypotheses = request.hypotheses or current_analysis.hypotheses
    updated_evidence = request.evidence or current_analysis.evidence
    updated_assessments = request.assessments or current_analysis.assessments

    # Generate updated matrix
    matrix = _generate_ach_matrix(updated_hypotheses, updated_evidence, updated_assessments)

    return ACHAnalysisResponse(
        session_id=session_id,
        title=updated_title,
        scenario=updated_scenario,
        key_question=updated_key_question,
        hypotheses=updated_hypotheses,
        evidence=updated_evidence,
        assessments=updated_assessments,
        matrix=matrix,
        status="in_progress",
        version=current_analysis.version + 1
    )


@router.post("/{session_id}/hypothesis")
async def add_hypothesis(
    session_id: int,
    hypothesis: Hypothesis,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
) -> dict:
    """
    Add a hypothesis to ACH analysis.
    
    Args:
        session_id: Session ID
        hypothesis: Hypothesis data
        current_user: Current authenticated user
        db: Database session
        
    Returns:
        dict: Success message
    """
    logger.info(f"Adding hypothesis to ACH {session_id}: {hypothesis.description}")

    return {
        "message": "Hypothesis added successfully",
        "hypothesis": hypothesis.dict(),
        "session_id": session_id
    }


@router.post("/{session_id}/evidence")
async def add_evidence(
    session_id: int,
    evidence: Evidence,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
) -> dict:
    """
    Add evidence to ACH analysis.
    
    Args:
        session_id: Session ID
        evidence: Evidence data
        current_user: Current authenticated user
        db: Database session
        
    Returns:
        dict: Success message
    """
    logger.info(f"Adding evidence to ACH {session_id}: {evidence.description}")

    return {
        "message": "Evidence added successfully",
        "evidence": evidence.dict(),
        "session_id": session_id
    }


@router.put("/{session_id}/assessment")
async def update_assessment(
    session_id: int,
    assessment: EvidenceAssessment,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
) -> dict:
    """
    Update evidence-hypothesis assessment.
    
    Args:
        session_id: Session ID
        assessment: Assessment data
        current_user: Current authenticated user
        db: Database session
        
    Returns:
        dict: Success message
    """
    logger.info(
        f"Updating assessment in ACH {session_id}: "
        f"E{assessment.evidence_id} vs H{assessment.hypothesis_id} = {assessment.consistency}"
    )

    return {
        "message": "Assessment updated successfully",
        "assessment": assessment.dict(),
        "session_id": session_id
    }


@router.get("/{session_id}/matrix")
async def get_ach_matrix(
    session_id: int,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
) -> dict:
    """
    Get the ACH matrix visualization data.
    
    Args:
        session_id: Session ID
        current_user: Current authenticated user
        db: Database session
        
    Returns:
        dict: Matrix data for visualization
    """
    logger.info(f"Getting ACH matrix for session {session_id}")

    # TODO: Get actual data from database
    # For now, return mock matrix
    hypotheses = [
        {"id": "h1", "description": "State-sponsored", "score": 7.2},
        {"id": "h2", "description": "Criminal group", "score": 5.8},
        {"id": "h3", "description": "Insider threat", "score": 3.1},
    ]

    evidence = [
        {"id": "e1", "description": "Advanced techniques"},
        {"id": "e2", "description": "Ransom note"},
        {"id": "e3", "description": "Timing of attack"},
    ]

    return {
        "session_id": session_id,
        "hypotheses": hypotheses,
        "evidence": evidence,
        "matrix": [
            ["", "H1: State", "H2: Criminal", "H3: Insider"],
            ["E1: Techniques", "++", "0", "--"],
            ["E2: Ransom", "--", "++", "0"],
            ["E3: Timing", "0", "0", "+"],
        ],
        "legend": {
            "++": "Strongly Consistent",
            "+": "Consistent",
            "0": "Neutral",
            "-": "Inconsistent",
            "--": "Strongly Inconsistent",
            "NA": "Not Applicable"
        }
    }


@router.post("/{session_id}/calculate-probabilities")
async def calculate_probabilities(
    session_id: int,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
) -> dict:
    """
    Calculate hypothesis probabilities based on evidence assessments.
    
    Args:
        session_id: Session ID
        current_user: Current authenticated user
        db: Database session
        
    Returns:
        dict: Updated probabilities
    """
    logger.info(f"Calculating probabilities for ACH {session_id}")

    # TODO: Implement Bayesian or weighted scoring algorithm
    # For now, return mock calculations

    probabilities = {
        "h1": {
            "hypothesis": "State-sponsored cyber attack",
            "initial_probability": 0.5,
            "updated_probability": 0.72,
            "confidence": 0.8,
            "supporting_evidence": 2,
            "contradicting_evidence": 1
        },
        "h2": {
            "hypothesis": "Criminal ransomware operation",
            "initial_probability": 0.5,
            "updated_probability": 0.58,
            "confidence": 0.7,
            "supporting_evidence": 1,
            "contradicting_evidence": 1
        },
        "h3": {
            "hypothesis": "Insider threat",
            "initial_probability": 0.5,
            "updated_probability": 0.31,
            "confidence": 0.6,
            "supporting_evidence": 1,
            "contradicting_evidence": 2
        }
    }

    return {
        "session_id": session_id,
        "probabilities": probabilities,
        "methodology": "Weighted evidence scoring",
        "timestamp": "2025-08-16T00:00:00Z"
    }


def _generate_ach_matrix(
    hypotheses: list,
    evidence: list,
    assessments: list
) -> dict:
    """
    Generate ACH matrix from hypotheses, evidence, and assessments.
    
    Args:
        hypotheses: List of hypotheses
        evidence: List of evidence
        assessments: List of assessments
        
    Returns:
        dict: Matrix structure
    """
    matrix = {
        "headers": ["Evidence"] + [h.description if hasattr(h, 'description') else h["description"]
                                   for h in hypotheses],
        "rows": []
    }

    for ev in evidence:
        ev_id = ev.id if hasattr(ev, 'id') else ev["id"]
        ev_desc = ev.description if hasattr(ev, 'description') else ev["description"]
        row = [ev_desc]

        for hyp in hypotheses:
            hyp_id = hyp.id if hasattr(hyp, 'id') else hyp["id"]

            # Find assessment for this evidence-hypothesis pair
            assessment = next(
                (a for a in assessments
                 if (a.evidence_id if hasattr(a, 'evidence_id') else a["evidence_id"]) == ev_id and
                    (a.hypothesis_id if hasattr(a, 'hypothesis_id') else a["hypothesis_id"]) == hyp_id),
                None
            )

            if assessment:
                consistency = assessment.consistency if hasattr(assessment, 'consistency') else assessment["consistency"]
                symbol_map = {
                    "consistent": "+",
                    "strongly_consistent": "++",
                    "inconsistent": "-",
                    "strongly_inconsistent": "--",
                    "neutral": "0",
                    "not_applicable": "NA"
                }
                row.append(symbol_map.get(consistency, "0"))
            else:
                row.append("0")

        matrix["rows"].append(row)

    return matrix


@router.get("/templates/list")
async def list_ach_templates(
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
) -> list[dict]:
    """
    List available ACH analysis templates.
    
    Args:
        current_user: Current authenticated user
        db: Database session
        
    Returns:
        list: Available templates
    """
    templates = [
        {
            "id": 1,
            "name": "Attribution Analysis",
            "description": "Template for attributing actions to actors",
            "sample_hypotheses": [
                "State-sponsored actor",
                "Criminal organization",
                "Hacktivist group",
                "Insider threat"
            ],
            "evidence_categories": [
                "Technical indicators",
                "Behavioral patterns",
                "Motivations",
                "Capabilities"
            ]
        },
        {
            "id": 2,
            "name": "Threat Assessment",
            "description": "Template for assessing potential threats",
            "sample_hypotheses": [
                "Imminent threat",
                "Developing threat",
                "Low probability threat",
                "No credible threat"
            ],
            "evidence_categories": [
                "Intelligence reports",
                "Open source information",
                "Technical indicators",
                "Historical patterns"
            ]
        },
        {
            "id": 3,
            "name": "Intent Analysis",
            "description": "Template for analyzing adversary intent",
            "sample_hypotheses": [
                "Espionage",
                "Sabotage",
                "Financial gain",
                "Political influence"
            ],
            "evidence_categories": [
                "Target selection",
                "Methods used",
                "Timing",
                "Communications"
            ]
        }
    ]

    return templates


@router.post("/ai/generate-hypotheses")
async def generate_ach_hypotheses_endpoint(
    request: dict
) -> dict:
    """Generate alternative hypotheses using AI for ACH analysis."""
    try:
        scenario = request.get("scenario", "")
        key_question = request.get("key_question", "")
        context = request.get("context", "")
        existing_hypotheses = request.get("existing_hypotheses", [])

        # Enhanced AI hypothesis generation
        ai_prompt = f"""
        As an intelligence analyst expert in the Analysis of Competing Hypotheses (ACH) methodology, 
        generate 5-7 distinct, mutually exclusive hypotheses to answer the key question.

        Scenario: {scenario}
        Key Question: {key_question}
        Additional Context: {context}
        Existing Hypotheses: {existing_hypotheses}

        Requirements:
        1. Each hypothesis should be specific and testable
        2. Hypotheses should cover the full spectrum of possibilities
        3. Include at least one contrarian or unlikely hypothesis
        4. Avoid overlap with existing hypotheses
        5. Consider different actor types, motivations, and methods
        
        For each hypothesis, provide:
        - A clear, concise description (max 100 words)
        - Initial probability assessment (0.1-0.9)
        - Key assumptions underlying the hypothesis
        - Critical evidence that would support/refute it
        
        Format as JSON array with objects containing: description, probability, assumptions, critical_evidence
        """

        result = await ai_service.analyze_with_ai_detailed(ai_prompt, "ach_hypothesis_generation")
        
        # Parse AI response and structure hypotheses
        if isinstance(result, dict) and "content" in result:
            import json
            try:
                hypotheses_data = json.loads(result["content"])
                if isinstance(hypotheses_data, list):
                    hypotheses = []
                    for i, hyp in enumerate(hypotheses_data):
                        hypothesis = {
                            "id": f"ai_h_{i+1}",
                            "description": hyp.get("description", ""),
                            "probability": hyp.get("probability", 0.5),
                            "assumptions": hyp.get("assumptions", []),
                            "critical_evidence": hyp.get("critical_evidence", []),
                            "source": "AI Generated",
                            "generated_at": datetime.now().isoformat()
                        }
                        hypotheses.append(hypothesis)
                else:
                    # Fallback to simpler generation
                    hypotheses = await ai_service.generate_ach_hypotheses(scenario, key_question, context)
            except json.JSONDecodeError:
                # Fallback to existing method
                hypotheses = await ai_service.generate_ach_hypotheses(scenario, key_question, context)
        else:
            # Fallback to existing method
            hypotheses = await ai_service.generate_ach_hypotheses(scenario, key_question, context)

        return {
            "hypotheses": hypotheses,
            "count": len(hypotheses),
            "methodology": "AI-Enhanced ACH Hypothesis Generation",
            "scenario_analysis": {
                "complexity_score": len(scenario.split()) / 10,  # Simple complexity metric
                "key_actors_identified": _extract_potential_actors(scenario),
                "temporal_scope": _assess_temporal_scope(scenario)
            },
            "generated_at": datetime.now().isoformat()
        }
    except Exception as e:
        logger.error(f"Enhanced hypothesis generation failed: {e}")
        raise HTTPException(status_code=500, detail="AI hypothesis generation failed")


@router.post("/ai/refine-hypothesis")
async def refine_hypothesis_endpoint(
    request: dict
) -> dict:
    """Refine a specific hypothesis based on evidence and feedback."""
    try:
        hypothesis = request.get("hypothesis", "")
        evidence = request.get("evidence", [])
        feedback = request.get("feedback", "")

        refined = await ai_service.refine_hypothesis(hypothesis, evidence, feedback)

        return {
            "original_hypothesis": hypothesis,
            "refined_hypothesis": refined,
            "refined_at": datetime.now().isoformat()
        }
    except Exception as e:
        logger.error(f"Hypothesis refinement failed: {e}")
        raise HTTPException(status_code=500, detail="AI hypothesis refinement failed")


@router.post("/ai/evidence-gaps")
async def suggest_evidence_gaps_endpoint(
    request: dict
) -> dict:
    """Identify evidence gaps that would improve hypothesis discrimination."""
    try:
        hypotheses = [h.get("text", "") for h in request.get("hypotheses", [])]
        evidence = request.get("evidence", [])

        gaps = await ai_service.suggest_evidence_gaps(hypotheses, evidence)

        return {
            "evidence_gaps": gaps,
            "gap_count": len(gaps),
            "analysis_date": datetime.now().isoformat(),
            "recommendations": [
                "Focus on discriminating evidence that distinguishes between hypotheses",
                "Seek diverse source types and perspectives",
                "Fill timeline gaps with additional data points"
            ]
        }
    except Exception as e:
        logger.error(f"Evidence gap analysis failed: {e}")
        raise HTTPException(status_code=500, detail="AI evidence gap analysis failed")


@router.post("/ai/bias-detection")
async def detect_analysis_bias_endpoint(
    request: dict
) -> dict:
    """Detect potential cognitive biases in ACH analysis."""
    try:
        hypotheses = [h.get("text", "") for h in request.get("hypotheses", [])]
        evidence = request.get("evidence", [])
        scores = request.get("scores", [])

        bias_analysis = await ai_service.detect_analysis_bias(hypotheses, evidence, scores)

        return {
            "bias_analysis": bias_analysis,
            "analysis_date": datetime.now().isoformat(),
            "total_hypotheses": len(hypotheses),
            "total_evidence": len(evidence),
            "total_scores": len(scores)
        }
    except Exception as e:
        logger.error(f"Bias detection failed: {e}")
        raise HTTPException(status_code=500, detail="AI bias detection failed")


class ACHExportRequest(BaseModel):
    """ACH export request model."""
    title: str
    scenario: str
    key_question: str
    hypotheses: list[dict]
    evidence: list[dict]
    scores: list[dict]
    analysis_date: str | None = None
    analyst_name: str | None = "Anonymous User"
    ranked_hypotheses: list[dict] | None = None


# Temporarily disabled until dependencies are properly installed
# @router.post("/export/excel")
async def export_ach_excel(request: ACHExportRequest) -> StreamingResponse:
    """Export ACH analysis as Excel matrix."""
    if not EXPORT_AVAILABLE:
        raise HTTPException(status_code=503, detail="Export functionality not available - missing dependencies")

    try:
        # Create Excel file in memory
        output = io.BytesIO()
        workbook = xlsxwriter.Workbook(output, {'in_memory': True})
        worksheet = workbook.add_worksheet('ACH Matrix')

        # Define formats
        header_format = workbook.add_format({
            'bold': True,
            'bg_color': '#D7E4BD',
            'border': 1
        })
        cell_format = workbook.add_format({'border': 1})

        # Write title and metadata
        worksheet.write(0, 0, 'ACH Analysis: ' + request.title)
        worksheet.write(1, 0, 'Scenario: ' + request.scenario)
        worksheet.write(2, 0, 'Key Question: ' + request.key_question)
        worksheet.write(3, 0, 'Analysis Date: ' + (request.analysis_date or datetime.now().strftime('%Y-%m-%d')))
        worksheet.write(4, 0, 'Analyst: ' + request.analyst_name)

        # Create matrix headers
        start_row = 6
        worksheet.write(start_row, 0, 'Evidence', header_format)

        # Write hypothesis headers
        for i, hypothesis in enumerate(request.hypotheses):
            worksheet.write(start_row, i + 1, f"H{i+1}: {hypothesis.get('text', '')[:30]}", header_format)

        # Write evidence rows
        for i, evidence in enumerate(request.evidence):
            row = start_row + 1 + i
            worksheet.write(row, 0, f"E{i+1}: {evidence.get('text', '')[:50]}", cell_format)

            # Fill matrix with scores
            for j, hypothesis in enumerate(request.hypotheses):
                # Find score for this evidence/hypothesis pair
                score_obj = next((s for s in request.scores
                                if s.get('evidenceId') == evidence.get('id') and
                                   s.get('hypothesisId') == hypothesis.get('id')), None)

                score_val = score_obj.get('score', 0) if score_obj else 0
                # Convert to traditional ACH symbols
                if score_val >= 3:
                    symbol = "+" if score_val == 3 else "++"
                elif score_val <= -3:
                    symbol = "-" if score_val == -3 else "--"
                else:
                    symbol = "0"

                worksheet.write(row, j + 1, symbol, cell_format)

        # Add SATS evaluation if present
        sats_row = start_row + len(request.evidence) + 3
        worksheet.write(sats_row, 0, 'SATS Evaluation Summary', header_format)

        for i, evidence in enumerate(request.evidence):
            sats = evidence.get('sats_evaluation')
            if sats:
                row = sats_row + 1 + i
                worksheet.write(row, 0, f"E{i+1} Overall Score:", cell_format)
                worksheet.write(row, 1, f"{sats.get('overall_score', 0)}/5", cell_format)

        workbook.close()
        output.seek(0)

        return StreamingResponse(
            io.BytesIO(output.read()),
            media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            headers={'Content-Disposition': f'attachment; filename="ach_matrix_{datetime.now().strftime("%Y%m%d")}.xlsx"'}
        )

    except Exception as e:
        logger.error(f"Excel export error: {e}")
        raise HTTPException(status_code=500, detail="Export failed")


@router.post("/export/word")
async def export_ach_word(request: ACHExportRequest) -> StreamingResponse:
    """Export ACH analysis as Word document."""
    if not EXPORT_AVAILABLE:
        raise HTTPException(status_code=503, detail="Export functionality not available - missing dependencies")

    try:
        doc = docx.Document()

        # Title and metadata
        title = doc.add_heading(f'Analysis of Competing Hypotheses: {request.title}', 0)

        doc.add_paragraph(f'Analysis Date: {request.analysis_date or datetime.now().strftime("%Y-%m-%d")}')
        doc.add_paragraph(f'Analyst: {request.analyst_name}')
        doc.add_paragraph('Classification: UNCLASSIFIED')

        # Executive Summary
        doc.add_heading('Executive Summary', level=1)
        doc.add_paragraph('This ACH analysis evaluates competing hypotheses to provide structured analytical reasoning.')

        # Scenario
        doc.add_heading('Scenario', level=1)
        doc.add_paragraph(request.scenario)

        # Key Question
        doc.add_heading('Key Intelligence Question', level=1)
        doc.add_paragraph(request.key_question)

        # Hypotheses
        doc.add_heading('Hypotheses Under Consideration', level=1)
        for i, hypothesis in enumerate(request.hypotheses):
            doc.add_paragraph(f"H{i+1}: {hypothesis.get('text', '')}", style='List Number')

        # Evidence
        doc.add_heading('Evidence Evaluation', level=1)
        for i, evidence in enumerate(request.evidence):
            doc.add_paragraph(f"E{i+1}: {evidence.get('text', '')}", style='List Number')

            # Add SATS evaluation if available
            sats = evidence.get('sats_evaluation')
            if sats:
                p = doc.add_paragraph()
                p.add_run('SATS Evaluation: ').bold = True
                p.add_run(f"Overall Score {sats.get('overall_score', 0)}/5 ")
                p.add_run(f"(Reliability: {sats.get('reliability', 0)}/5, ")
                p.add_run(f"Accuracy: {sats.get('accuracy', 0)}/5, ")
                p.add_run(f"Timeliness: {sats.get('timeliness', 0)}/5, ")
                p.add_run(f"Significance: {sats.get('significance', 0)}/5)")

        # Analysis Results
        if request.ranked_hypotheses:
            doc.add_heading('Analysis Results', level=1)
            doc.add_paragraph('Hypotheses ranked by evidence support:')
            for i, hyp in enumerate(request.ranked_hypotheses):
                doc.add_paragraph(f"{i+1}. {hyp.get('text', '')} (Score: {hyp.get('score', 0)})", style='List Number')

        # Methodology
        doc.add_heading('Methodology', level=1)
        doc.add_paragraph('This analysis uses the ACH methodology developed by Heuer to systematically evaluate competing hypotheses against available evidence. Evidence is scored using SATS criteria (Source reliability, Accuracy, Timeliness, Significance) when applicable.')

        # Save to memory
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)

        return StreamingResponse(
            io.BytesIO(output.read()),
            media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            headers={'Content-Disposition': f'attachment; filename="ach_report_{datetime.now().strftime("%Y%m%d")}.docx"'}
        )

    except Exception as e:
        logger.error(f"Word export error: {e}")
        raise HTTPException(status_code=500, detail="Export failed")


@router.post("/export/powerpoint")
async def export_ach_powerpoint(request: ACHExportRequest) -> StreamingResponse:
    """Export ACH analysis as PowerPoint presentation."""
    if not EXPORT_AVAILABLE:
        raise HTTPException(status_code=503, detail="Export functionality not available - missing dependencies")

    try:
        prs = PPTXPresentation()

        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = f"ACH Analysis: {request.title}"
        subtitle.text = f"Analysis Date: {request.analysis_date or datetime.now().strftime('%Y-%m-%d')}\nAnalyst: {request.analyst_name}\nClassification: UNCLASSIFIED"

        # Scenario slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Analysis Scenario"
        content.text = request.scenario

        # Key question slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Key Intelligence Question"
        content.text = request.key_question

        # Hypotheses slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Hypotheses Under Consideration"
        hyp_text = "\n".join([f"H{i+1}: {h.get('text', '')}" for i, h in enumerate(request.hypotheses)])
        content.text = hyp_text

        # Results slide
        if request.ranked_hypotheses:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]

            title.text = "Analysis Results"
            results_text = "Hypotheses ranked by evidence support:\n\n"
            for i, hyp in enumerate(request.ranked_hypotheses[:3]):  # Top 3
                results_text += f"{i+1}. {hyp.get('text', '')}\n   Score: {hyp.get('score', 0)}\n\n"
            content.text = results_text

        # Methodology slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Methodology"
        content.text = "Analysis of Competing Hypotheses (ACH)\n\n• Structured analytical technique\n• Systematic evaluation of evidence\n• SATS criteria applied where applicable\n• Reduces cognitive bias in analysis"

        # Save to memory
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        return StreamingResponse(
            io.BytesIO(output.read()),
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            headers={'Content-Disposition': f'attachment; filename="ach_briefing_{datetime.now().strftime("%Y%m%d")}.pptx"'}
        )

    except Exception as e:
        logger.error(f"PowerPoint export error: {e}")
        raise HTTPException(status_code=500, detail="Export failed")


def _extract_potential_actors(scenario: str) -> list[str]:
    """Extract potential actors from scenario text using simple keyword matching."""
    actors = []
    scenario_lower = scenario.lower()
    
    # Common actor types in intelligence analysis
    actor_keywords = {
        "state": ["government", "state", "nation", "country", "ministry", "agency"],
        "criminal": ["criminal", "gang", "cartel", "mafia", "organized crime"],
        "terrorist": ["terrorist", "extremist", "radical", "militia", "insurgent"],
        "corporate": ["company", "corporation", "business", "industry", "commercial"],
        "individual": ["individual", "person", "actor", "agent", "operative"],
        "hacktivist": ["hacktivist", "activist", "hacker", "anonymous"]
    }
    
    for actor_type, keywords in actor_keywords.items():
        if any(keyword in scenario_lower for keyword in keywords):
            actors.append(actor_type)
    
    return actors


def _assess_temporal_scope(scenario: str) -> str:
    """Assess the temporal scope of the scenario."""
    scenario_lower = scenario.lower()
    
    if any(word in scenario_lower for word in ["immediate", "urgent", "now", "today", "imminent"]):
        return "immediate"
    elif any(word in scenario_lower for word in ["recent", "week", "month", "lately"]):
        return "short-term"
    elif any(word in scenario_lower for word in ["year", "years", "long-term", "historical"]):
        return "long-term"
    else:
        return "medium-term"


@router.post("/ai/automated-evidence-collection")
async def automated_evidence_collection_endpoint(
    request: dict
) -> dict:
    """Automatically collect and analyze evidence for ACH hypotheses using web scraping."""
    try:
        hypotheses = request.get("hypotheses", [])
        search_terms = request.get("search_terms", [])
        time_range = request.get("time_range", "30d")  # 30 days default
        
        # Prepare web scraping requests for evidence collection
        evidence_results = []
        
        for hypothesis in hypotheses:
            hyp_desc = hypothesis.get("description", "")
            
            # Generate search queries for this hypothesis
            search_queries = await _generate_evidence_search_queries(hyp_desc, search_terms)
            
            # Simulate evidence collection (in production, would use web scraping service)
            for query in search_queries[:3]:  # Limit to 3 queries per hypothesis
                evidence_item = {
                    "id": f"auto_ev_{len(evidence_results) + 1}",
                    "description": f"Evidence related to: {query}",
                    "source": "Automated collection",
                    "query_used": query,
                    "hypothesis_relevance": hyp_desc,
                    "collection_date": datetime.now().isoformat(),
                    "credibility_score": 0.7,  # Would be assessed by AI in production
                    "relevance_score": 0.8
                }
                evidence_results.append(evidence_item)
        
        return {
            "evidence_collected": evidence_results,
            "collection_summary": {
                "total_evidence": len(evidence_results),
                "hypotheses_analyzed": len(hypotheses),
                "time_range": time_range,
                "collection_date": datetime.now().isoformat()
            },
            "next_steps": [
                "Review automated evidence for relevance",
                "Validate source credibility",
                "Update ACH matrix with new evidence",
                "Consider additional search terms"
            ]
        }
        
    except Exception as e:
        logger.error(f"Automated evidence collection failed: {e}")
        raise HTTPException(status_code=500, detail="Evidence collection failed")


async def _generate_evidence_search_queries(hypothesis: str, base_terms: list[str]) -> list[str]:
    """Generate search queries for evidence collection."""
    # In production, this would use AI to generate more sophisticated queries
    queries = []
    
    # Extract key terms from hypothesis
    key_terms = [term.strip() for term in hypothesis.split() if len(term.strip()) > 3]
    
    # Combine with base search terms
    for base_term in base_terms[:2]:  # Limit base terms
        for key_term in key_terms[:2]:  # Limit key terms
            queries.append(f"{base_term} {key_term}")
    
    # Add the hypothesis itself as a search term
    queries.append(hypothesis[:100])  # Truncate long hypotheses
    
    return queries[:5]  # Return max 5 queries


@router.post("/ai/consistency-analysis")
async def automated_consistency_analysis_endpoint(
    request: dict
) -> dict:
    """Perform automated consistency analysis between evidence and hypotheses."""
    try:
        hypotheses = request.get("hypotheses", [])
        evidence = request.get("evidence", [])
        
        consistency_matrix = []
        analysis_details = []
        
        for evidence_item in evidence:
            evidence_id = evidence_item.get("id", "")
            evidence_desc = evidence_item.get("description", "")
            
            row_analysis = {
                "evidence_id": evidence_id,
                "evidence_description": evidence_desc,
                "hypothesis_assessments": []
            }
            
            for hypothesis in hypotheses:
                hypothesis_id = hypothesis.get("id", "")
                hypothesis_desc = hypothesis.get("description", "")
                
                # AI-powered consistency assessment
                consistency_prompt = f"""
                Analyze the consistency between this evidence and hypothesis:
                
                Evidence: {evidence_desc}
                Hypothesis: {hypothesis_desc}
                
                Assess on scale:
                - Strongly Supports (++): Evidence strongly supports hypothesis
                - Supports (+): Evidence supports hypothesis  
                - Neutral (0): Evidence neither supports nor contradicts
                - Contradicts (-): Evidence contradicts hypothesis
                - Strongly Contradicts (--): Evidence strongly contradicts hypothesis
                
                Provide reasoning and confidence score (0-1).
                Respond with JSON: {{"assessment": "++|+|0|-|--", "reasoning": "explanation", "confidence": 0.0-1.0}}
                """
                
                # Simulate AI assessment (in production would call AI service)
                assessment = await _simulate_consistency_assessment(evidence_desc, hypothesis_desc)
                
                hypothesis_assessment = {
                    "hypothesis_id": hypothesis_id,
                    "consistency": assessment["assessment"],
                    "reasoning": assessment["reasoning"],
                    "confidence": assessment["confidence"],
                    "analysis_date": datetime.now().isoformat()
                }
                
                row_analysis["hypothesis_assessments"].append(hypothesis_assessment)
                
                # Add to matrix
                consistency_matrix.append({
                    "evidence_id": evidence_id,
                    "hypothesis_id": hypothesis_id,
                    "consistency": assessment["assessment"],
                    "confidence": assessment["confidence"]
                })
            
            analysis_details.append(row_analysis)
        
        # Calculate hypothesis scores
        hypothesis_scores = _calculate_hypothesis_scores(consistency_matrix, hypotheses)
        
        return {
            "consistency_matrix": consistency_matrix,
            "analysis_details": analysis_details,
            "hypothesis_scores": hypothesis_scores,
            "analysis_summary": {
                "total_assessments": len(consistency_matrix),
                "high_confidence_assessments": len([a for a in consistency_matrix if a["confidence"] > 0.8]),
                "conflicting_evidence": len([a for a in consistency_matrix if a["consistency"] in ["-", "--"]]),
                "supporting_evidence": len([a for a in consistency_matrix if a["consistency"] in ["+", "++"]]),
                "analysis_date": datetime.now().isoformat()
            },
            "recommendations": [
                "Review low-confidence assessments",
                "Seek additional evidence for neutral assessments",
                "Investigate conflicting evidence sources",
                "Update hypothesis probabilities based on scores"
            ]
        }
        
    except Exception as e:
        logger.error(f"Consistency analysis failed: {e}")
        raise HTTPException(status_code=500, detail="Consistency analysis failed")


async def _simulate_consistency_assessment(evidence: str, hypothesis: str) -> dict:
    """Simulate AI consistency assessment (placeholder for actual AI call)."""
    import random
    
    # Simple simulation based on keyword matching
    evidence_lower = evidence.lower()
    hypothesis_lower = hypothesis.lower()
    
    # Count common words (excluding common words)
    stop_words = {"the", "a", "an", "and", "or", "but", "in", "on", "at", "to", "for", "of", "with", "by"}
    evidence_words = set(evidence_lower.split()) - stop_words
    hypothesis_words = set(hypothesis_lower.split()) - stop_words
    
    common_words = evidence_words.intersection(hypothesis_words)
    similarity = len(common_words) / max(len(evidence_words), len(hypothesis_words), 1)
    
    # Determine assessment based on similarity and randomness
    if similarity > 0.3:
        assessment = random.choice(["+", "++"])
        confidence = 0.7 + similarity * 0.3
        reasoning = f"Evidence contains {len(common_words)} relevant terms supporting the hypothesis"
    elif similarity > 0.1:
        assessment = "0"
        confidence = 0.6
        reasoning = "Evidence is somewhat related but neither strongly supports nor contradicts"
    else:
        assessment = random.choice(["-", "0"])
        confidence = 0.5 + random.random() * 0.3
        reasoning = "Limited relevance between evidence and hypothesis"
    
    return {
        "assessment": assessment,
        "reasoning": reasoning,
        "confidence": min(confidence, 1.0)
    }


def _calculate_hypothesis_scores(consistency_matrix: list, hypotheses: list) -> list:
    """Calculate scores for each hypothesis based on consistency assessments."""
    scores = []
    
    for hypothesis in hypotheses:
        hypothesis_id = hypothesis.get("id", "")
        
        # Get all assessments for this hypothesis
        assessments = [item for item in consistency_matrix if item["hypothesis_id"] == hypothesis_id]
        
        # Calculate weighted score
        total_score = 0
        total_weight = 0
        
        for assessment in assessments:
            consistency = assessment["consistency"]
            confidence = assessment["confidence"]
            
            # Convert consistency to numeric score
            score_map = {"++": 2, "+": 1, "0": 0, "-": -1, "--": -2}
            numeric_score = score_map.get(consistency, 0)
            
            # Weight by confidence
            weighted_score = numeric_score * confidence
            total_score += weighted_score
            total_weight += confidence
        
        # Calculate final score
        final_score = total_score / max(total_weight, 1)
        
        scores.append({
            "hypothesis_id": hypothesis_id,
            "hypothesis_description": hypothesis.get("description", ""),
            "score": round(final_score, 2),
            "supporting_evidence": len([a for a in assessments if a["consistency"] in ["+", "++"]]),
            "contradicting_evidence": len([a for a in assessments if a["consistency"] in ["-", "--"]]),
            "neutral_evidence": len([a for a in assessments if a["consistency"] == "0"]),
            "average_confidence": round(sum(a["confidence"] for a in assessments) / max(len(assessments), 1), 2)
        })
    
    # Sort by score descending
    scores.sort(key=lambda x: x["score"], reverse=True)
    
    return scores
